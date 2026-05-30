"""
JCC Assistant - Bible Study & Church Programs Chatbot
======================================================
Jubilee Celebration Center - AFM.
"""

import os
import json
import base64
from datetime import date
from pathlib import Path
from typing import List, Tuple, Optional

import gradio as gr
from supabase import create_client, Client
from openai import OpenAI
from docx import Document
from pptx import Presentation


# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------
SUPABASE_URL = os.environ["SUPABASE_URL"]
SUPABASE_KEY = os.environ["SUPABASE_SERVICE_KEY"]
OPENAI_API_KEY = os.environ["OPENAI_API_KEY"]
ADMIN_PASSWORD = os.environ["ADMIN_PASSWORD"]

OPENAI_MODEL = "gpt-4o-mini"

sb: Client = create_client(SUPABASE_URL, SUPABASE_KEY)
oai = OpenAI(api_key=OPENAI_API_KEY)


# ---------------------------------------------------------------------------
# Logo as base64 (so it inlines into the header HTML)
# ---------------------------------------------------------------------------
LOGO_PATH = Path(__file__).parent / "jcc_logo.jpeg"
if LOGO_PATH.exists():
    with open(LOGO_PATH, "rb") as f:
        LOGO_B64 = base64.b64encode(f.read()).decode("ascii")
    LOGO_DATA_URI = f"data:image/jpeg;base64,{LOGO_B64}"
else:
    LOGO_DATA_URI = ""


# ---------------------------------------------------------------------------
# Document parsing
# ---------------------------------------------------------------------------
def parse_docx(file_path: str) -> str:
    doc = Document(file_path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def parse_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    sections = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_lines = [f"## Slide {i}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(run.text for run in para.runs).strip()
                    if txt:
                        slide_lines.append(txt)
        if len(slide_lines) > 1:
            sections.append("\n".join(slide_lines))
    return "\n\n".join(sections)


def parse_study_document(file_path: str) -> str:
    lower = file_path.lower()
    if lower.endswith(".docx"):
        return parse_docx(file_path)
    elif lower.endswith(".pptx"):
        return parse_pptx(file_path)
    else:
        raise ValueError("Unsupported file type. Please upload .docx or .pptx.")


# ---------------------------------------------------------------------------
# Suggested questions generator
# ---------------------------------------------------------------------------
SUGGEST_PROMPT = """Below is a Bible study document. Generate exactly 5 short, specific questions a group member might ask about THIS study after reading it.

Rules:
- Each question is 4-10 words.
- Questions must be answerable from the document content.
- Cover different sections / angles (not all the same topic).
- No generic Bible questions like "What does the Bible say about faith?" - they must be specific to this study.

Return a JSON array of 5 strings only. No preamble, no markdown, no code fences. Example:
["What does the study teach about X?", "What scriptures support Y?", ...]

STUDY DOCUMENT:
{document_text}
"""


def generate_suggested_questions(document_text: str) -> list:
    """Generate 5 study-specific questions. Returns [] on failure."""
    try:
        resp = oai.chat.completions.create(
            model=OPENAI_MODEL,
            messages=[{"role": "user", "content": SUGGEST_PROMPT.format(document_text=document_text)}],
            temperature=0.5,
        )
        raw = resp.choices[0].message.content.strip()
        # strip code fences if model added them
        if raw.startswith("```"):
            raw = raw.split("```")[1]
            if raw.startswith("json"):
                raw = raw[4:]
            raw = raw.strip()
        questions = json.loads(raw)
        if isinstance(questions, list):
            return [str(q) for q in questions[:5]]
    except Exception as e:
        print(f"Question generation failed: {e}")
    return []


PROGRAMS_SUGGESTED = [
    "When is the next Couples Ministry event?",
    "What is the vision for Praise & Worship?",
    "What fundraising activities are planned for 2026?",
    "Who leads the Outreach Ministry?",
    "What events are happening in May 2026?",
]


# ---------------------------------------------------------------------------
# Bible Study DB helpers
# ---------------------------------------------------------------------------
def list_bible_studies() -> List[Tuple[str, str]]:
    res = sb.table("bible_studies") \
        .select("id, week_of, title, presenter") \
        .order("week_of", desc=True) \
        .order("uploaded_at", desc=True) \
        .execute()
    options = []
    for row in res.data:
        label = f"{row['week_of']} - {row['title']}"
        if row.get("presenter"):
            label += f" ({row['presenter']})"
        options.append((label, row["id"]))
    return options


def get_bible_study(study_id: str) -> Optional[dict]:
    if not study_id:
        return None
    res = sb.table("bible_studies").select("*").eq("id", study_id).single().execute()
    return res.data


def upload_bible_study(file, title, presenter, week_of, password):
    if password != ADMIN_PASSWORD:
        return "❌ Incorrect admin password."
    if not file:
        return "❌ Please attach a .docx or .pptx file."
    if not title or not title.strip():
        return "❌ Title is required."
    if not week_of:
        return "❌ Week-of date is required."

    try:
        text = parse_study_document(file.name)
    except Exception as e:
        return f"❌ Could not parse the document: {e}"

    if len(text) < 50:
        return "❌ The document looks empty after parsing."

    # Generate 5 study-specific questions
    questions = generate_suggested_questions(text)

    # Check for existing study with same week_of -> replace
    existing = sb.table("bible_studies").select("id").eq("week_of", week_of).execute()
    payload = {
        "title": title.strip(),
        "presenter": (presenter or "").strip() or None,
        "document_text": text,
        "suggested_questions": questions,
    }

    if existing.data:
        existing_id = existing.data[0]["id"]
        try:
            sb.table("bible_studies").update(payload).eq("id", existing_id).execute()
            return (
                f"✅ Replaced existing study for week of {week_of} with **{title}**.\n\n"
                f"Document length: {len(text):,} characters. "
                f"Generated {len(questions)} suggested questions."
            )
        except Exception as e:
            return f"❌ Database update failed: {e}"

    payload["week_of"] = week_of
    try:
        sb.table("bible_studies").insert(payload).execute()
    except Exception as e:
        return f"❌ Database insert failed: {e}"

    return (
        f"✅ Uploaded **{title}** for the week of {week_of}.\n\n"
        f"Document length: {len(text):,} characters. "
        f"Generated {len(questions)} suggested questions."
    )


def delete_bible_study(study_id, password):
    if password != ADMIN_PASSWORD:
        return "❌ Incorrect admin password.", gr.update()
    if not study_id:
        return "❌ Please select a study to delete.", gr.update()
    try:
        sb.table("bible_studies").delete().eq("id", study_id).execute()
    except Exception as e:
        return f"❌ Delete failed: {e}", gr.update()
    options = list_bible_studies()
    return "✅ Deleted.", gr.update(choices=options, value=options[0][1] if options else None)


# ---------------------------------------------------------------------------
# Church Programs context
# ---------------------------------------------------------------------------
def fetch_programs_context() -> str:
    ministries = sb.table("ministries").select("*").execute().data
    events = sb.table("events").select("*").order("event_date").execute().data
    notes = sb.table("ministry_notes").select("*").execute().data

    lines = ["# JCC 2026 MINISTRY PROGRAMS\n"]
    for m in ministries:
        lines.append(f"\n## {m['name']} Ministry")
        if m.get("lead"):
            lines.append(f"Led by: {m['lead']}")

        m_notes = [n for n in notes if n["ministry_id"] == m["id"]]
        for n in m_notes:
            heading = n["section"].replace("_", " ").title()
            lines.append(f"\n**{heading}:** {n['content']}")

        m_events = [e for e in events if e["ministry_id"] == m["id"]]
        if m_events:
            lines.append("\n**Events:**")
            for e in m_events:
                date_str = e.get("date_label") or e.get("event_date") or "TBD"
                line = f"- {date_str}: {e['title']}"
                if e.get("description"):
                    line += f" - {e['description']}"
                if e.get("format"):
                    line += f" ({e['format']})"
                lines.append(line)
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# System prompts
# ---------------------------------------------------------------------------
BIBLE_STUDY_PROMPT = """You are the JCC Bible Study Assistant for Jubilee Celebration Center - AFM.

You answer questions about a specific Bible study document shared with the group. Your answers MUST come only from the document below - do not add interpretation, outside scripture, or commentary.

RULES:
- Quote directly from the document using quotation marks for the key passage.
- Cite the section heading, number, or slide where the quote appears (e.g., "Section 7: Logos and Rhema" or "Slide 4").
- If the question is not addressed in the document, say:
  "This isn't covered in this week's study. Please bring it up with the group or the presenter."
- Do not speculate or fill gaps with general Bible knowledge.
- Keep answers focused. Concise is better than thorough.
- Maintain a respectful, reverent tone appropriate for Bible study.

---
THIS WEEK'S BIBLE STUDY
Title: {title}
Presenter: {presenter}
Week of: {week_of}

{document_text}
---
"""

PROGRAMS_PROMPT = """You are the JCC Programs Assistant for Jubilee Celebration Center - AFM.

You answer questions about JCC's 2026 ministry programs, events, leads, goals, and activities. Your answers MUST come only from the programs information below.

GUIDANCE:
- Be specific. When asked about events, give the date, ministry, and format.
- When asked about a ministry's vision, mission, or goals, summarize from the notes.
- Some questions are about people. Look across ALL ministries in the data - leads are listed at each ministry section. If a person appears as a lead of any ministry, mention that.
- For example, if asked "Who is the Pastor?", note that Pastor Tabu Bere leads the Outreach Ministry, even though there is no separate "Pastor" entry.
- If the information truly is not in the data below, say:
  "I don't have that information in the 2026 plans. Please check with the relevant ministry lead."
- Do not invent events, dates, or leads. Stick to the data.
- Keep answers focused and helpful.

---
JCC 2026 PROGRAMS DATA
{programs_context}
---
"""


# ---------------------------------------------------------------------------
# Chat
# ---------------------------------------------------------------------------
def chat(message, history, mode, study_id):
    if not message or not message.strip():
        return "Please type a question."

    if mode == "Bible Study":
        study = get_bible_study(study_id) if study_id else None
        if not study:
            return (
                "Please select a Bible study from the dropdown first. "
                "If the dropdown is empty, click Refresh or ask an admin to upload one."
            )
        system_prompt = BIBLE_STUDY_PROMPT.format(
            title=study["title"],
            presenter=study.get("presenter") or "Not specified",
            week_of=study["week_of"],
            document_text=study["document_text"],
        )
    else:
        try:
            context = fetch_programs_context()
        except Exception as e:
            return f"Could not load programs data: {e}"
        system_prompt = PROGRAMS_PROMPT.format(programs_context=context)

    messages = [{"role": "system", "content": system_prompt}]
    for h in history:
        messages.append({"role": h["role"], "content": h["content"]})
    messages.append({"role": "user", "content": message})

    try:
        resp = oai.chat.completions.create(
            model=OPENAI_MODEL,
            messages=messages,
            temperature=0.2,
        )
        return resp.choices[0].message.content
    except Exception as e:
        return f"OpenAI request failed: {e}"


# ---------------------------------------------------------------------------
# Dynamic suggested questions for sidebar
# ---------------------------------------------------------------------------
def get_suggestions(mode: str, study_id: str) -> List[str]:
    if mode == "Church Programs":
        return PROGRAMS_SUGGESTED
    # Bible Study
    if not study_id:
        return ["(Select a study to see suggestions)"]
    try:
        study = get_bible_study(study_id)
        if study and study.get("suggested_questions"):
            qs = study["suggested_questions"]
            if isinstance(qs, list) and qs:
                return qs
        return [
            "What were the main points?",
            "What scriptures are referenced?",
            "Summarize the conclusion",
            "What does the study say about the key topic?",
            "What is the application for our lives?",
        ]
    except Exception as e:
        print(f"get_suggestions failed: {e}")
        return []


def refresh_studies_and_suggestions(mode):
    options = list_bible_studies()
    if not options:
        suggestions = get_suggestions(mode, None)
        return (
            gr.update(choices=[], value=None),
            *[gr.update(value=q, visible=True) for q in suggestions[:5]],
            *[gr.update(visible=False) for _ in range(5 - len(suggestions))],
        )
    new_value = options[0][1]
    suggestions = get_suggestions(mode, new_value)
    visible_count = min(len(suggestions), 5)
    updates = []
    for i in range(5):
        if i < visible_count:
            updates.append(gr.update(value=suggestions[i], visible=True))
        else:
            updates.append(gr.update(visible=False))
    return (gr.update(choices=options, value=new_value), *updates)


def update_suggestions(mode, study_id):
    suggestions = get_suggestions(mode, study_id)
    updates = []
    for i in range(5):
        if i < len(suggestions):
            updates.append(gr.update(value=suggestions[i], visible=True))
        else:
            updates.append(gr.update(visible=False))
    return updates


# ---------------------------------------------------------------------------
# UI
# ---------------------------------------------------------------------------
CUSTOM_CSS = """
.gradio-container {
    font-family: 'Inter', 'Helvetica Neue', system-ui, sans-serif !important;
    max-width: 1280px !important;
    margin: 0 auto !important;
}
#jcc-hero {
    background: linear-gradient(135deg, #1B2A4E 0%, #2C4170 100%);
    border-radius: 16px;
    padding: 28px 32px;
    margin-bottom: 18px;
    color: white;
    display: flex;
    align-items: center;
    gap: 24px;
    box-shadow: 0 8px 24px rgba(27, 42, 78, 0.18);
    position: relative;
    overflow: hidden;
}
#jcc-hero::after {
    content: "";
    position: absolute;
    bottom: 0; left: 0; right: 0;
    height: 4px;
    background: linear-gradient(90deg, #C9A55C 0%, #E4CC8E 50%, #C9A55C 100%);
}
#jcc-hero img.logo {
    width: 88px;
    height: 88px;
    border-radius: 50%;
    background: white;
    padding: 6px;
    box-shadow: 0 4px 12px rgba(0,0,0,0.2);
    flex-shrink: 0;
}
#jcc-hero .titles h1 {
    font-size: 1.9em !important;
    font-weight: 700 !important;
    margin: 0 0 4px 0 !important;
    color: white !important;
    letter-spacing: -0.5px;
}
#jcc-hero .titles .church-name {
    font-size: 0.85em;
    color: #C9A55C;
    letter-spacing: 3px;
    font-weight: 600;
    margin-bottom: 6px;
    text-transform: uppercase;
}
#jcc-hero .titles .tagline {
    font-size: 0.95em;
    color: #cbd5e1;
    margin: 0;
}
.sidebar-card {
    background: white;
    border-radius: 12px;
    padding: 16px;
    border: 1px solid #e5e7eb;
    margin-bottom: 12px;
}
.sidebar-card h3 {
    color: #1B2A4E;
    font-size: 0.78em;
    letter-spacing: 1.5px;
    text-transform: uppercase;
    margin: 0 0 12px 0;
    font-weight: 700;
    border-left: 3px solid #C9A55C;
    padding-left: 10px;
}
.suggest-btn button {
    background: white !important;
    border: 1px solid #e5e7eb !important;
    color: #1B2A4E !important;
    text-align: left !important;
    font-weight: 500 !important;
    font-size: 0.88em !important;
    padding: 10px 12px !important;
    line-height: 1.35 !important;
    white-space: normal !important;
    height: auto !important;
    min-height: 40px !important;
    transition: all 0.15s ease;
    width: 100% !important;
    justify-content: flex-start !important;
}
.suggest-btn button:hover {
    background: #1B2A4E !important;
    color: white !important;
    border-color: #1B2A4E !important;
    transform: translateX(2px);
}
.tab-nav button {
    font-weight: 500 !important;
}
.tab-nav button.selected {
    color: #1B2A4E !important;
    border-bottom-color: #C9A55C !important;
}
footer { display: none !important; }
"""


theme = gr.themes.Soft(
    primary_hue=gr.themes.colors.slate,
    secondary_hue=gr.themes.colors.amber,
    neutral_hue=gr.themes.colors.slate,
    font=[gr.themes.GoogleFont("Inter"), "system-ui", "sans-serif"],
).set(
    button_primary_background_fill="#1B2A4E",
    button_primary_background_fill_hover="#0F1A35",
    button_primary_text_color="white",
    body_background_fill="#F7F3EC",
    block_background_fill="white",
    block_border_color="#e5e7eb",
)


with gr.Blocks(title="JCC Assistant", theme=theme, css=CUSTOM_CSS) as demo:

    # ============================================================
    # HERO HEADER with logo
    # ============================================================
    logo_img_html = (
        f'<img class="logo" src="{LOGO_DATA_URI}" alt="JCC Logo"/>'
        if LOGO_DATA_URI else ""
    )
    gr.HTML(f"""
    <div id="jcc-hero">
        {logo_img_html}
        <div class="titles">
            <div class="church-name">Jubilee Celebration Center &mdash; AFM</div>
            <h1>JCC Assistant</h1>
            <p class="tagline">Ask about this week's Bible study or church programs and events for 2026.</p>
        </div>
    </div>
    """)

    with gr.Tabs():
        # ============================================================
        # CHAT TAB
        # ============================================================
        with gr.Tab("Chat"):
            with gr.Row():
                # LEFT SIDEBAR
                with gr.Column(scale=1, min_width=260):
                    with gr.Group(elem_classes=["sidebar-card"]):
                        gr.HTML("<h3>Mode</h3>")
                        mode = gr.Radio(
                            choices=["Bible Study", "Church Programs"],
                            value="Bible Study",
                            show_label=False,
                            container=False,
                        )

                    with gr.Group(elem_classes=["sidebar-card"]):
                        gr.HTML("<h3>Bible Study</h3>")
                        study_dropdown = gr.Dropdown(
                            choices=list_bible_studies(),
                            show_label=False,
                            container=False,
                        )
                        refresh_btn = gr.Button("Refresh list", size="sm")

                    with gr.Group(elem_classes=["sidebar-card"]):
                        gr.HTML("<h3>Suggested Questions</h3>")
                        suggest_btns = [
                            gr.Button("", elem_classes=["suggest-btn"], visible=False)
                            for _ in range(5)
                        ]

                # MAIN CHAT
                with gr.Column(scale=3):
                    chatbot = gr.Chatbot(
                        type="messages",
                        height=560,
                        avatar_images=(None, LOGO_DATA_URI or None),
                        show_label=False,
                        bubble_full_width=False,
                    )
                    msg = gr.Textbox(
                        placeholder="Ask a question…",
                        show_label=False,
                        container=False,
                        autofocus=True,
                    )

            # ---------- chat plumbing ----------
            def respond(message, history, mode_val, study_id):
                if not message or not message.strip():
                    return "", history
                history = history + [{"role": "user", "content": message}]
                reply = chat(message, history[:-1], mode_val, study_id)
                history = history + [{"role": "assistant", "content": reply}]
                return "", history

            msg.submit(
                respond,
                inputs=[msg, chatbot, mode, study_dropdown],
                outputs=[msg, chatbot],
            )

            def send_suggested(question, history, mode_val, study_id):
                if not question or not question.strip():
                    return history
                history = history + [{"role": "user", "content": question}]
                reply = chat(question, history[:-1], mode_val, study_id)
                history = history + [{"role": "assistant", "content": reply}]
                return history

            for btn in suggest_btns:
                btn.click(
                    fn=send_suggested,
                    inputs=[btn, chatbot, mode, study_dropdown],
                    outputs=[chatbot],
                )

            # When mode or study changes, refresh suggestions
            mode.change(
                fn=update_suggestions,
                inputs=[mode, study_dropdown],
                outputs=suggest_btns,
            )
            study_dropdown.change(
                fn=update_suggestions,
                inputs=[mode, study_dropdown],
                outputs=suggest_btns,
            )
            refresh_btn.click(
                fn=refresh_studies_and_suggestions,
                inputs=[mode],
                outputs=[study_dropdown, *suggest_btns],
            )

            # Init suggestions on load
            demo.load(
                fn=update_suggestions,
                inputs=[mode, study_dropdown],
                outputs=suggest_btns,
            )

        # ============================================================
        # ADMIN TAB
        # ============================================================
        with gr.Tab("Admin"):
            gr.Markdown(
                "### Upload a Bible Study\n"
                "Accepts `.docx` or `.pptx` files. "
                "If a study already exists for the same week-of date, it will be **replaced**. "
                "Five suggested questions are auto-generated for each upload."
            )

            with gr.Row():
                with gr.Column(scale=2):
                    up_title = gr.Textbox(
                        label="Title",
                        placeholder="e.g. The Foundation of the Word",
                    )
                    up_presenter = gr.Textbox(
                        label="Presenter",
                        placeholder="e.g. Elder, Group 3",
                    )
                    up_week = gr.Textbox(
                        label="Week of (YYYY-MM-DD)",
                        value=str(date.today()),
                    )
                with gr.Column(scale=1):
                    up_password = gr.Textbox(
                        label="Admin Password",
                        type="password",
                    )
                    up_file = gr.File(
                        label="Document (.docx / .pptx)",
                        file_types=[".docx", ".pptx"],
                    )

            up_button = gr.Button("Upload Study", variant="primary")
            up_status = gr.Markdown()

            up_button.click(
                fn=upload_bible_study,
                inputs=[up_file, up_title, up_presenter, up_week, up_password],
                outputs=up_status,
            )

            gr.Markdown("---\n### Delete a Bible Study")
            with gr.Row():
                del_dropdown = gr.Dropdown(
                    label="Select study to delete",
                    choices=list_bible_studies(),
                    scale=4,
                )
                del_refresh = gr.Button("Refresh list", scale=1, size="sm")
                del_password = gr.Textbox(
                    label="Admin Password",
                    type="password",
                    scale=2,
                )
            del_button = gr.Button("Delete Selected", variant="stop")
            del_status = gr.Markdown()

            def refresh_del():
                options = list_bible_studies()
                return gr.update(choices=options, value=options[0][1] if options else None)

            del_refresh.click(fn=refresh_del, outputs=del_dropdown)
            del_button.click(
                fn=delete_bible_study,
                inputs=[del_dropdown, del_password],
                outputs=[del_status, del_dropdown],
            )

    gr.HTML(
        "<div style='text-align:center; color:#9ca3af; font-size:0.85em; padding:12px;'>"
        "JCC Assistant - Prototype. The bot only answers from loaded study and programs data."
        "</div>"
    )


if __name__ == "__main__":
    demo.launch()
